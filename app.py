import g4f # This allows us to write text using gpt-4
from pptx import Presentation # Creates Presentation for us
import random # For generating random numbers
import re # For regular expressions 
import os # Allow us to interact with system
import speech_recognition as sr # recognize speech
import pyttsx3 # A module for tts (text-to-speech)
import time

engine = pyttsx3.init("sapi5")
voices = engine.getProperty("voices")
engine.setProperty("voice", voices[2].id)
engine.setProperty("rate", 155)

def speak(audio):
    engine.say(audio)
    engine.runAndWait()

def takecommand():

    r = sr.Recognizer()

    with sr.Microphone() as source:
        print('listening....')
        r.pause_threshold = 1
        r.adjust_for_ambient_noise(source)
        
        audio = r.listen(source, 10, 6)

    try:
        print("Understanding..")
        query  = r.recognize_google(audio,language='en-in')
        print(f"You Said: {query}\n")
        time.sleep(1)
       
    except Exception as e:
        print(f"error: {e}")
        return "None"
    return query

Prompt_template = """Create a powerpoint/presentation about the user's topic. You only answer with presentation. Follow the structure of the example.
Notice
-You only do all the presentation text for user
-You write the text no longer than 250 character!
-Make short titles.
-Make it easy to understand
-The presentation has a table of content
-The presentation has summary
-At least 7 slides
-The end slide should be a thank you slide

Example! - Stick to this formatting exactly!
#TITLE: TITLE OF THE PRESENTATION

#SLIDE: 1
#HEADER: Table of content
#CONTENT: 1. CONTENT OF THIS PRESENTATION
2. CONTENT OF THIS PRESENTATION
3. CONTENT OF THIS PRESENTATION

#SLIDE: 2
#HEADER: TITLE OF THIS SLIDE
#CONTENT: CONTENT OF THIS SLIDE

#SLIDE: 3
#HEADER: TITLE OF THIS SLIDE
#CONTENT: CONTENT OF THIS SLIDE

#SLIDE: 4
#HEADER: TITLE OF THIS SLIDE
#CONTENT: CONTENT OF THIS SLIDE

#SLIDE: 5
#HEADER: SUMMARY
#CONTENT: CONTENT THE OF SUMMARY

#SLIDE: END
#CONTENT: Thank you"""

# Function for generating presentation text
def generate_presentation_text(user_input):
    # Call the gpt 4 api for generating text based upon the user's input
    response = g4f.ChatCompletion.create(
        model="llama-3-70b", # Specify The gpt model here
        # model="gpt-4-32k-0163", # Specify The gpt model here
        provider=g4f.Provider.MetaAI,
        messages=[
            {"role": "system", "content": Prompt_template}, #Prompt for generating presentation
            {"role": "user", "content": f"User wants a presentation about {user_input}"}, # User role for input
        ],
        stream=True, #Enable streaming for real-time response
    )

    # Extract and concatenate the generated text from gpt-4 response
    generated_text = ""
    for message in response:
        generated_text += str(message)
        # Print each message for real-time feedback during generation process
        print(message, end="", flush=True)
    print()
     
    return generated_text # Giving generated text as final output

# A function to create powerpoint presentation
def create_presentation(text_file, design_number, presentation_name):
    # Ensure the output directory exists
    output_directory = 'GeneratedPresentation'
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)
        
    # Load the selected design layout
    presentation = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0 # initializing the slide count
    slide_title = "" # initializing the slide Title
    slide_content = "" # initializing the slide content
    last_slide_layout_index = -1 # initializing the last_slide_layout_index
    first_time = True # Flag for first time entering the loop
    # Read the text file line by line to structure the presentation
    with open(text_file, 'r', encoding='utf-8') as file:
        for line_num, line in enumerate(file):
            # Detect title slide
            if line.startswith("#TITLE: "):
                slide_title = line.replace("#TITLE: ","").strip()
                slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                title_shapes = slide.shapes.title
                title_shapes.text = slide_title
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith("#SLIDE: "):
                if slide_count > 0:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[slide_layout_index])
                    title_shape = slide.shapes.title
                    title_shape.text = slide_title
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    text_frame = body_shape.text_frame
                    text_frame.text = slide_content

                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1,7,8]
                while slide_layout_index == last_slide_layout_index:
                    if first_time:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        first_time = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                            
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith("#HEADER: "):
                slide_title = line.replace("#HEADER: ","").strip()
                continue
            elif line.startswith("#CONTENT: "):
                slide_content = line.replace("#CONTENT: ","").strip()
                next_line = file.readline().strip()
                while next_line and not next_line.startswith("#"):
                    slide_content += '\n' + next_line
                    next_line = file.readline().strip()
                continue
        
        # Save the presentation
        presentation.save(f'{output_directory}/{presentation_name}.pptx')
        file_path = f"{output_directory}/{presentation_name}.pptx"

        return file_path
    

def get_presentation(user_input, startfile=True):
    user_text = user_input 
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]','',input_string)
    input_string = input_string.replace("\n","")
    design_number = user_input # Default design number
    speak("Processing...")

    if last_char.isdigit():
        design_number = int(last_char)
        input_string = user_text[0:-2]
        print("Design Number", design_number, "selected")

    else:
        print("No design specified, using default design...")

    if design_number > "7" or design_number == "0":
        design_number = 2
        print("Design number unavailable, using default design...")

    with open(f'cache/{input_string}.txt', 'w', encoding='utf-8') as file:
        file.write(generate_presentation_text(input_string))

    presentation_link = create_presentation(f'cache/{input_string}.txt', design_number, input_string)
    return str(presentation_link)

speak("Done Sir")
